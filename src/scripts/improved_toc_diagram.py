#!/usr/bin/env python3
"""
Improved Theory of Change Diagram Generator
This script creates a single PowerPoint slide with an enhanced, professional
Theory of Change diagram for resource sharing, learning, and training.
"""

import os
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.path import Path
import numpy as np
from PIL import Image

# Define colors that match the Clarity Impact Finance brand
PRIMARY_COLOR = "#1B4620"  # Dark Green
SECONDARY_COLOR = "#E67E45"  # Orange
ACCENT_COLOR_1 = "#153618"  # Darker Green
ACCENT_COLOR_2 = "#d16a33"  # Darker Orange
LIGHT_GREEN = "#3E8C48"  # Light Green
MED_GREEN = "#2D6934"  # Medium Green
BACKGROUND_COLOR = "#FFFFFF"  # White

def create_improved_toc_diagram():
    """Create an improved Theory of Change diagram using matplotlib."""
    # Create a figure with a white background
    fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
    plt.tight_layout(rect=[0, 0, 1, 0.95])  # Adjust the layout to leave room for title
    
    # Turn off axis
    ax.axis('off')
    ax.set_xlim(-0.5, 4.5)
    ax.set_ylim(-4.5, 2.5)  # Extend the bottom limit to make room for text
    
    # Define the components of the Theory of Change
    components = {
        'inputs': ['Knowledge Repository', 'Expert Network', 'Learning Platform', 'Community Engagement'],
        'activities': ['Resource Curation', 'Collaborative Learning', 'Skill-Building Workshops', 'Peer Exchange'],
        'outputs': ['Practical Toolkits', 'Case Studies', 'Training Modules', 'Community of Practice'],
        'outcomes': ['Enhanced CDFI Capacity', 'Improved Lending Practices', 'Stronger Compliance', 'Operational Efficiency'],
        'impact': ['Increased Community Investment', 'Greater CDFI Sustainability', 'Expanded Access to Capital', 'Stronger Communities']
    }
    
    # Define positions for each component
    positions = {
        'inputs': 0,
        'activities': 1,
        'outputs': 2, 
        'outcomes': 3,
        'impact': 4
    }
    
    # Define colors for each component
    colors = {
        'inputs': PRIMARY_COLOR,
        'activities': MED_GREEN,
        'outputs': LIGHT_GREEN,
        'outcomes': SECONDARY_COLOR,
        'impact': ACCENT_COLOR_2
    }
    
    # Add a title
    plt.suptitle('Theory of Change: Resources, Learning, and Training Framework', 
                fontsize=24, fontweight='bold', color=PRIMARY_COLOR, y=0.98)
    
    # Draw the main flow arrows first (behind the boxes)
    arrow_y = 0.5
    for i in range(len(positions)-1):
        component_names = list(positions.keys())
        start_x = positions[component_names[i]] + 0.4
        end_x = positions[component_names[i+1]] - 0.4
        
        # Draw a curved arrow
        arrow = patches.FancyArrowPatch(
            (start_x, arrow_y), (end_x, arrow_y),
            connectionstyle="arc3,rad=0.1",
            arrowstyle="fancy,head_width=0.4,head_length=0.6,tail_width=0.2",
            fc='gray', ec='gray', alpha=0.8,
            linewidth=2
        )
        ax.add_patch(arrow)
    
    # Draw component boxes and their items
    box_width = 0.8
    box_height = 1.0
    
    # First, draw all main component boxes
    for component, items in components.items():
        x = positions[component]
        
        # Draw the main component box with rounded corners
        rect = patches.FancyBboxPatch(
            (x - box_width/2, arrow_y - box_height/2),
            box_width, box_height,
            boxstyle=patches.BoxStyle("Round", pad=0.1),
            facecolor=colors[component], alpha=0.9,
            edgecolor='black', linewidth=1.5
        )
        ax.add_patch(rect)
        
        # Add component title
        plt.text(x, arrow_y,
                component.upper(),
                ha='center', va='center',
                color='white', fontweight='bold', fontsize=16,
                bbox=dict(facecolor='none', edgecolor='none', pad=10))
    
    # Now add item boxes with more space between them
    for component, items in components.items():
        x = positions[component]
        
        # Add items below the component box with more spacing
        for i, item in enumerate(items):
            y_pos = -1.0 - i*0.8  # Increased spacing between items
            
            # Create a rounded rectangle with stronger border for each item
            item_rect = patches.FancyBboxPatch(
                (x - 0.9, y_pos - 0.25),
                1.8, 0.5,  # Slightly larger boxes
                boxstyle=patches.BoxStyle("Round", pad=0.2),
                facecolor='white', alpha=0.95,
                edgecolor=colors[component], linewidth=2.0
            )
            ax.add_patch(item_rect)
            
            # Add a subtle shadow effect
            shadow_rect = patches.FancyBboxPatch(
                (x - 0.9 + 0.03, y_pos - 0.25 + 0.03),
                1.8, 0.5,
                boxstyle=patches.BoxStyle("Round", pad=0.2),
                facecolor='gray', alpha=0.1,
                edgecolor=None
            )
            ax.add_patch(shadow_rect)
            
            # Add text with improved visibility
            plt.text(x, y_pos,
                    item,
                    ha='center', va='center',
                    color=colors[component], fontsize=13, fontweight='bold',
                    bbox=dict(facecolor='none', edgecolor='none', pad=10))
    
    # Add connecting lines from main boxes to their items
    for component, items in components.items():
        x = positions[component]
        
        # Draw a vertical line connecting the main box to its items
        if items:
            # Start from the bottom of the main box
            start_y = arrow_y - box_height/2
            # End at the top of the first item box
            end_y = -1.0 + 0.25
            
            line = plt.Line2D([x, x], [start_y, end_y], 
                             color=colors[component], linewidth=1.5, 
                             alpha=0.7, linestyle='--')
            ax.add_line(line)
    
    # Add explanatory text
    plt.figtext(0.5, 0.02, 
               "This framework illustrates how our resources and activities lead to meaningful impact for CDFIs and communities.",
               ha='center', fontsize=14, color='#333333', weight='bold')
    
    # Save the figure to a bytes buffer
    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=300, bbox_inches='tight')
    buf.seek(0)
    
    return buf

def create_single_toc_slide():
    """Create a single PowerPoint slide with the improved Theory of Change diagram."""
    prs = Presentation()
    
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Add a blank slide
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Create the diagram
    diagram_buffer = create_improved_toc_diagram()
    
    # Add the diagram to the slide
    img = Image.open(diagram_buffer)
    img_bytes = io.BytesIO()
    img.save(img_bytes, format='PNG')
    img_bytes.seek(0)
    
    # Add picture to slide, centered and sized appropriately
    slide.shapes.add_picture(img_bytes, Inches(0.5), Inches(0.25), 
                           width=Inches(12.33), height=Inches(7.0))
    
    # Create the scripts directory if it doesn't exist
    os.makedirs('src/scripts', exist_ok=True)
    
    # Save the presentation
    output_path = 'src/scripts/improved_toc_slide.pptx'
    prs.save(output_path)
    print(f"Improved Theory of Change slide saved to {output_path}")

if __name__ == "__main__":
    create_single_toc_slide()
    plt.close('all')  # Close all matplotlib figures 