#!/usr/bin/env python3
"""
Theory of Change Framework PowerPoint Generator
This script creates a PowerPoint presentation that illustrates a theory of change
framework for resource sharing, learning, and training.
"""

import os
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import numpy as np
from PIL import Image

# Define colors that match the Clarity Impact Finance brand
PRIMARY_COLOR = RGBColor(27, 70, 32)  # #1B4620 - Dark Green
SECONDARY_COLOR = RGBColor(230, 126, 69)  # #E67E45 - Orange
ACCENT_COLOR_1 = RGBColor(21, 54, 24)  # #153618 - Darker Green
ACCENT_COLOR_2 = RGBColor(209, 106, 51)  # #d16a33 - Darker Orange
BACKGROUND_COLOR = RGBColor(255, 255, 255)  # White

def create_title_slide(prs):
    """Create the title slide for the presentation."""
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Resources Hub: Theory of Change Framework"
    subtitle.text = "A Framework for Learning, Resource Sharing, and Training"
    
    # Customize title formatting
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Customize subtitle formatting
    subtitle.text_frame.paragraphs[0].font.color.rgb = SECONDARY_COLOR
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    
    return slide

def create_introduction_slide(prs):
    """Create an introduction slide explaining the purpose of the framework."""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Introduction to Our Theory of Change"
    
    # Add content with bullet points
    text_frame = content.text_frame
    text_frame.clear()  # Clear existing content
    
    p = text_frame.paragraphs[0]
    p.text = "Our Theory of Change framework guides how we approach resource sharing, learning, and training for CDFIs:"
    p.font.size = Pt(18)
    p.font.color.rgb = PRIMARY_COLOR
    
    bullet_points = [
        "We believe that accessible, practical knowledge leads to stronger CDFIs",
        "Learning should be collaborative and community-driven",
        "Resources must be actionable and tailored to CDFI needs",
        "Training should build both technical skills and strategic thinking",
        "Impact measurement is integrated throughout the learning journey"
    ]
    
    for point in bullet_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(16)
    
    return slide

def create_toc_diagram(prs):
    """Create a slide with the Theory of Change diagram."""
    # Create a matplotlib figure for the diagram
    plt.figure(figsize=(10, 6))
    plt.tight_layout()
    
    # Turn off axis
    ax = plt.gca()
    ax.axis('off')
    
    # Define the components of the Theory of Change
    components = {
        'inputs': ['Knowledge Repository', 'Expert Network', 'Learning Platform', 'Community Engagement'],
        'activities': ['Resource Curation', 'Collaborative Learning', 'Skill-Building Workshops', 'Peer Exchange'],
        'outputs': ['Practical Toolkits', 'Case Studies', 'Training Modules', 'Community of Practice'],
        'outcomes': ['Enhanced CDFI Capacity', 'Improved Lending Practices', 'Stronger Compliance', 'Operational Efficiency'],
        'impact': ['Increased Community Investment', 'Greater CDFI Sustainability', 'Expanded Access to Capital', 'Stronger Communities']
    }
    
    # Define positions for each component
    positions = {
        'inputs': 0,
        'activities': 1,
        'outputs': 2, 
        'outcomes': 3,
        'impact': 4
    }
    
    # Define colors for each component
    colors = {
        'inputs': '#1B4620',      # Dark Green
        'activities': '#2D6934',  # Medium Green
        'outputs': '#3E8C48',     # Light Green
        'outcomes': '#E67E45',    # Orange
        'impact': '#d16a33'       # Dark Orange
    }
    
    # Draw the boxes and arrows
    box_height = 0.7
    arrow_props = dict(arrowstyle='->', connectionstyle='arc3,rad=0.1', 
                      color='gray', lw=2)
    
    # Draw component boxes
    for component, items in components.items():
        x = positions[component]
        
        # Draw the main component box
        rect = plt.Rectangle((x-0.4, -0.5), 0.8, box_height, 
                           color=colors[component], alpha=0.8, 
                           ec='black', lw=1)
        ax.add_patch(rect)
        
        # Add component title
        plt.text(x, -0.15, component.upper(), 
                ha='center', va='center', 
                color='white', fontweight='bold', fontsize=12)
        
        # Add items below the component box
        for i, item in enumerate(items):
            y_pos = 0.4 + i*0.3
            plt.text(x, y_pos, item, ha='center', va='center', 
                    color=colors[component], fontsize=10, 
                    bbox=dict(facecolor='white', alpha=0.8, boxstyle='round,pad=0.5'))
    
    # Draw arrows between components
    for i in range(len(positions)-1):
        component_names = list(positions.keys())
        start = positions[component_names[i]]
        end = positions[component_names[i+1]]
        plt.annotate('', xy=(end-0.4, 0), xytext=(start+0.4, 0), arrowprops=arrow_props)
    
    # Add title
    plt.title('Theory of Change: Resources, Learning, and Training Framework', 
             fontsize=16, fontweight='bold', color='#1B4620', y=1.05)
    
    # Save the figure to a bytes buffer
    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=300, bbox_inches='tight')
    buf.seek(0)
    
    # Add a slide for the diagram
    slide_layout = prs.slide_layouts[5]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Theory of Change Framework"
    
    # Add the diagram to the slide
    img = Image.open(buf)
    # Convert PIL Image to bytes
    img_bytes = io.BytesIO()
    img.save(img_bytes, format='PNG')
    img_bytes.seek(0)
    
    # Add picture to slide
    slide.shapes.add_picture(img_bytes, Inches(0.5), Inches(1.5), 
                           width=Inches(9), height=Inches(5))
    
    plt.close()
    
    return slide

def create_components_slides(prs, components):
    """Create detailed slides for each component of the Theory of Change."""
    component_descriptions = {
        'inputs': "The foundational resources and assets we leverage to drive change in the CDFI ecosystem.",
        'activities': "The key actions and initiatives we undertake to transform inputs into valuable outputs.",
        'outputs': "The tangible products, services, and deliverables that result from our activities.",
        'outcomes': "The medium-term changes and benefits that occur as a result of our outputs.",
        'impact': "The long-term, sustainable changes in the CDFI ecosystem and communities served."
    }
    
    component_details = {
        'inputs': [
            "Knowledge Repository: Curated collection of best practices, research, and industry standards",
            "Expert Network: Practitioners, academics, and policy experts who contribute specialized knowledge",
            "Learning Platform: Digital infrastructure for hosting and delivering educational content",
            "Community Engagement: Relationships with CDFIs, funders, and community stakeholders"
        ],
        'activities': [
            "Resource Curation: Identifying, organizing, and contextualizing valuable information",
            "Collaborative Learning: Facilitating peer-to-peer knowledge exchange and co-creation",
            "Skill-Building Workshops: Delivering targeted training on technical and strategic topics",
            "Peer Exchange: Creating opportunities for CDFIs to share experiences and solutions"
        ],
        'outputs': [
            "Practical Toolkits: Ready-to-use templates, guides, and frameworks for CDFI operations",
            "Case Studies: Documented examples of successful approaches and lessons learned",
            "Training Modules: Structured learning experiences on key topics like underwriting and compliance",
            "Community of Practice: Active network of practitioners sharing knowledge and support"
        ],
        'outcomes': [
            "Enhanced CDFI Capacity: Strengthened organizational capabilities and staff skills",
            "Improved Lending Practices: More effective and equitable approaches to underwriting and lending",
            "Stronger Compliance: Better adherence to regulatory requirements and program guidelines",
            "Operational Efficiency: Streamlined processes and reduced administrative burden"
        ],
        'impact': [
            "Increased Community Investment: More capital flowing to underserved communities",
            "Greater CDFI Sustainability: Improved financial and operational performance of CDFIs",
            "Expanded Access to Capital: More individuals and businesses receiving needed financing",
            "Stronger Communities: Enhanced economic resilience and opportunity in target areas"
        ]
    }
    
    colors = {
        'inputs': RGBColor(27, 70, 32),      # Dark Green
        'activities': RGBColor(45, 105, 52),  # Medium Green
        'outputs': RGBColor(62, 140, 72),     # Light Green
        'outcomes': RGBColor(230, 126, 69),    # Orange
        'impact': RGBColor(209, 106, 51)       # Dark Orange
    }
    
    slides = []
    
    for component, items in components.items():
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        slides.append(slide)
        
        title = slide.shapes.title
        title.text = f"{component.title()}: {component_descriptions[component]}"
        
        # Format title
        title.text_frame.paragraphs[0].font.color.rgb = colors[component]
        
        # Add content placeholder
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        # Add component details
        for detail in component_details[component]:
            p = text_frame.add_paragraph()
            p.text = detail
            p.level = 0
            p.font.size = Pt(18)
            
            # Extract the title part (before the colon)
            if ":" in detail:
                title_part, description = detail.split(":", 1)
                p.text = title_part + ":"
                p.font.bold = True
                p.font.color.rgb = colors[component]
                
                # Add the description as a sub-bullet
                p2 = text_frame.add_paragraph()
                p2.text = description.strip()
                p2.level = 1
                p2.font.size = Pt(16)
    
    return slides

def create_implementation_slide(prs):
    """Create a slide about implementing the Theory of Change."""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Implementing the Framework"
    
    # Add content with bullet points
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "Key strategies for putting our Theory of Change into practice:"
    p.font.size = Pt(18)
    p.font.color.rgb = PRIMARY_COLOR
    
    implementation_points = [
        "Phased Approach: Begin with high-impact resources and gradually expand offerings",
        "Feedback Loops: Continuously gather input from CDFIs to refine resources",
        "Partnerships: Collaborate with industry associations, funders, and experts",
        "Technology: Leverage digital platforms for broader reach and accessibility",
        "Measurement: Track engagement, adoption, and impact metrics throughout"
    ]
    
    for point in implementation_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(16)
        
        # Make the first part bold
        if ":" in point:
            parts = point.split(":", 1)
            p.text = parts[0] + ":"
            p.font.bold = True
            
            # Add the rest as normal text
            p2 = text_frame.add_paragraph()
            p2.text = parts[1].strip()
            p2.level = 2
            p2.font.size = Pt(16)
    
    return slide

def create_conclusion_slide(prs):
    """Create a conclusion slide."""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Our Commitment to Learning and Growth"
    
    # Add content
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "Through our Resources Hub and Theory of Change framework, we commit to:"
    p.font.size = Pt(18)
    p.font.color.rgb = PRIMARY_COLOR
    
    commitments = [
        "Building a dynamic knowledge ecosystem that evolves with CDFI needs",
        "Fostering a culture of continuous learning and improvement",
        "Democratizing access to high-quality resources and expertise",
        "Measuring and sharing the impact of our learning initiatives",
        "Collaborating with partners to strengthen the entire CDFI ecosystem"
    ]
    
    for commitment in commitments:
        p = text_frame.add_paragraph()
        p.text = commitment
        p.level = 1
        p.font.size = Pt(16)
    
    # Add a call to action
    p = text_frame.add_paragraph()
    p.text = "\nJoin us in transforming how CDFIs learn, share, and grow together."
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = SECONDARY_COLOR
    p.font.bold = True
    
    return slide

def main():
    """Main function to create the PowerPoint presentation."""
    prs = Presentation()
    
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Create slides
    create_title_slide(prs)
    create_introduction_slide(prs)
    
    # Define the components of the Theory of Change
    components = {
        'inputs': ['Knowledge Repository', 'Expert Network', 'Learning Platform', 'Community Engagement'],
        'activities': ['Resource Curation', 'Collaborative Learning', 'Skill-Building Workshops', 'Peer Exchange'],
        'outputs': ['Practical Toolkits', 'Case Studies', 'Training Modules', 'Community of Practice'],
        'outcomes': ['Enhanced CDFI Capacity', 'Improved Lending Practices', 'Stronger Compliance', 'Operational Efficiency'],
        'impact': ['Increased Community Investment', 'Greater CDFI Sustainability', 'Expanded Access to Capital', 'Stronger Communities']
    }
    
    create_toc_diagram(prs)
    create_components_slides(prs, components)
    create_implementation_slide(prs)
    create_conclusion_slide(prs)
    
    # Create the scripts directory if it doesn't exist
    os.makedirs('src/scripts', exist_ok=True)
    
    # Save the presentation
    output_path = 'src/scripts/theory_of_change_framework.pptx'
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    main() 